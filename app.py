from flask import Flask, request, jsonify, send_file, render_template
from openai import OpenAI
from pydantic import BaseModel, ValidationError
from pptx import Presentation
from pptx.util import Inches, Pt
from weasyprint import HTML
import os
import io
import logging
import sqlite3
import json
from dotenv import load_dotenv
from database import init_db, save_prompt_template, get_prompt_templates

app = Flask(__name__)
load_dotenv()

# Configurar logging
logging.basicConfig(filename='app.log', level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Configurar la API de OpenAI
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

# Inicializar base de datos
init_db()

# Modelo para validar entrada
class PresentationRequest(BaseModel):
    topic: str
    num_slides: int
    tone: str
    audience: str
    language: str
    presentation_type: str
    theme: str
    context: str = ""

@app.route('/')
def index():
    try:
        templates = get_prompt_templates()
        return render_template('index.html', templates=templates)
    except Exception as e:
        logging.error(f"Error al cargar index.html: {str(e)}")
        return jsonify({"error": "No se pudo cargar la página. Verifica la configuración."}), 500

@app.route('/generate', methods=['POST'])
def generate_presentation():
    try:
        data = PresentationRequest(**request.get_json())
        if data.num_slides < 1 or data.num_slides > 10:
            raise ValueError("El número de diapositivas debe estar entre 1 y 10.")
        if data.tone not in ["formal", "informal"]:
            raise ValueError("El tono debe ser 'formal' o 'informal'.")
        if data.audience not in ["estudiantes", "profesionales"]:
            raise ValueError("El público debe ser 'estudiantes' o 'profesionales'.")
        if data.language not in ["es", "en", "fr"]:
            raise ValueError("El idioma debe ser 'es', 'en' o 'fr'.")
        if data.presentation_type not in ["Educativa", "Corporativa", "Creativa"]:
            raise ValueError("El tipo de presentación debe ser 'Educativa', 'Corporativa' o 'Creativa'.")
        if data.theme not in ["Minimalista", "Corporativo", "Colorido"]:
            raise ValueError("El tema debe ser 'Minimalista', 'Corporativo' o 'Colorido'.")

        # Calcular max_tokens dinámicamente
        max_tokens = 200 * data.num_slides + 500

        # Generar esquema previo
        outline_prompt = f"""
        Genera un esquema para una presentación {data.presentation_type} sobre "{data.topic}" en formato JSON, en idioma {data.language}.
        Incluye un título, un resumen ejecutivo (máximo 50 palabras) y {data.num_slides} secciones, cada una con un título atractivo (máximo 10 palabras).
        Considera el tono {data.tone}, el público {data.audience} y el contexto: {data.context}.
        Formato:
        {{
            "title": "Título de la presentación",
            "summary": "Resumen ejecutivo",
            "outline": [
                {{"title": "Título sección 1"}},
                ...
            ]
        }}
        """
        outline_response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": outline_prompt}],
            max_tokens=200
        )
        # Log raw response and token usage
        raw_outline = outline_response.choices[0].message.content.strip('```json\n').strip('```')
        logging.info(f"Raw outline response: {raw_outline}")
        logging.info(f"Outline token usage: {outline_response.usage.total_tokens}")
        # Parse JSON safely
        try:
            outline = json.loads(raw_outline)
        except json.JSONDecodeError as e:
            logging.error(f"Failed to parse outline JSON: {str(e)}")
            raise ValueError("Error al parsear el esquema de la presentación.")

        # Generar presentación completa
        prompt = f"""
        Genera una presentación {data.presentation_type} sobre "{data.topic}" en formato JSON, en idioma {data.language}.
        Usa el siguiente esquema: {outline}.
        Incluye un resumen ejecutivo como primera diapositiva.
        Cada diapositiva debe tener:
        - Un título atractivo (máximo 10 palabras).
        - Contenido (máximo 100 palabras).
        - Notas del orador (máximo 50 palabras).
        - Una descripción de imagen para ilustrar el contenido (máximo 20 palabras).
        Considera el tono {data.tone}, el público {data.audience} y el contexto: {data.context}.
        Asegura que el contenido sea legible para {data.audience} (nivel de lectura: {'8º grado' if data.audience == 'estudiantes' else 'universitario'}).
        Aplica un estilo visual {data.theme} (sugiere colores y fuentes).
        Incluye 2 preguntas abiertas para la audiencia al final y una cita relevante (máximo 30 palabras).
        Formato:
        {{
            "title": "Título de la presentación",
            "summary": "Resumen ejecutivo",
            "slides": [
                {{"title": "Resumen Ejecutivo", "content": "Contenido del resumen", "speaker_notes": "Notas del orador", "image_description": "Descripción de imagen"}},
                {{"title": "Título diapositiva 1", "content": "Contenido diapositiva 1", "speaker_notes": "Notas del orador", "image_description": "Descripción de imagen"}},
                ...
            ],
            "audience_questions": ["Pregunta 1", "Pregunta 2"],
            "quote": "Cita relevante",
            "theme_styles": {{"background": "color", "text_color": "color", "font": "fuente"}}
        }}
        """
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=max_tokens
        )
        # Log raw response and token usage
        raw_response = response.choices[0].message.content.strip('```json\n').strip('```')
        logging.info(f"Raw presentation response: {raw_response}")
        logging.info(f"Presentation token usage: {response.usage.total_tokens}")
        # Parse JSON safely
        try:
            presentation = json.loads(raw_response)
        except json.JSONDecodeError as e:
            logging.error(f"Failed to parse presentation JSON: {str(e)}")
            # Attempt to fix truncated JSON
            if raw_response.endswith('",') or raw_response.endswith('"'):
                try:
                    # Truncate the last incomplete string and close the JSON
                    last_comma = raw_response.rfind(',')
                    if last_comma != -1:
                        fixed_response = raw_response[:last_comma] + ']}]}'
                        presentation = json.loads(fixed_response)
                        logging.info("Fixed truncated JSON by completing the structure.")
                    else:
                        raise
                except json.JSONDecodeError:
                    raise ValueError("Error al parsear la presentación generada: JSON incompleto.")
            else:
                raise ValueError("Error al parsear la presentación generada.")

        # Generar imágenes con DALL·E
        for slide in presentation['slides']:
            if not slide.get('image_description'):
                slide['image_url'] = ''
                continue
            try:
                image_response = client.images.generate(
                    model="dall-e-2",
                    prompt=slide['image_description'],
                    size="512x512",
                    n=1
                )
                slide['image_url'] = image_response.data[0].url
            except Exception as e:
                logging.error(f"Error al generar imagen para diapositiva: {str(e)}")
                slide['image_url'] = ''

        # Moderación opcional (skip si falla)
        try:
            for slide in presentation['slides']:
                moderation = client.moderations.create(input=slide['content'] + slide['speaker_notes'])
                if moderation.results[0].flagged:
                    raise ValueError("Contenido inapropiado detectado en una diapositiva.")
        except Exception as e:
            logging.warning(f"Moderación falló: {str(e)}. Continuando sin moderación.")

        logging.info(f"Presentación generada para el tema: {data.topic}")
        return jsonify(presentation)
    except ValidationError as e:
        logging.error(f"Error de validación: {str(e)}")
        return jsonify({"error": "Datos de entrada no válidos."}), 400
    except Exception as e:
        logging.error(f"Error al generar presentación: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/export_pdf', methods=['POST'])
def export_pdf():
    try:
        data = request.get_json()
        html_content = render_template('slide_template.html', slides=data['slides'], title=data['title'], summary=data.get('summary', ''), audience_questions=data.get('audience_questions', []), quote=data.get('quote', ''))
        pdf_file = io.BytesIO()
        HTML(string=html_content).write_pdf(pdf_file)
        pdf_file.seek(0)
        return send_file(pdf_file, download_name='presentation.pdf', as_attachment=True)
    except Exception as e:
        logging.error(f"Error al exportar PDF: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/export_pptx', methods=['POST'])
def export_pptx():
    try:
        data = request.get_json()
        prs = Presentation()
        for slide_data in data['slides']:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = slide_data['title']
            content.text = slide_data['content'] + "\n\nNotas: " + slide_data['speaker_notes']
        pptx_file = io.BytesIO()
        prs.save(pptx_file)
        pptx_file.seek(0)
        return send_file(pptx_file, download_name='presentation.pptx', as_attachment=True)
    except Exception as e:
        logging.error(f"Error al exportar PPTX: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/export_json', methods=['POST'])
def export_json():
    try:
        data = request.get_json()
        json_file = io.StringIO()
        json_file.write(json.dumps(data))
        json_file.seek(0)
        return send_file(json_file, download_name='presentation.json', as_attachment=True)
    except Exception as e:
        logging.error(f"Error al exportar JSON: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/save_template', methods=['POST'])
def save_template():
    try:
        data = request.get_json()
        template_name = data['name']
        template_content = data['content']
        save_prompt_template(template_name, template_content)
        return jsonify({"message": "Plantilla guardada correctamente."})
    except Exception as e:
        logging.error(f"Error al guardar plantilla: {str(e)}")
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)